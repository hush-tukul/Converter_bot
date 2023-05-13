import asyncio
import io
import logging
import os
import subprocess

import aiofiles
import cv2
import numpy as np

from docx2txt import docx2txt

from pptx import Presentation
from docx2pdf import convert
from pdf2docx import parse
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter

from run import bot
from io import BytesIO

async def conv_to_jpeg(file_path: str, user_id) -> bytes:
    save_to_io = io.BytesIO()
    converted = fr"C:\PY\Python_learn\Minions_Bots\Converter_bot\tgbot_template_v3\save_dir\new_temp-{user_id}.jpg"   #add user_name to atribute!!!!!!!!!!!!!

    with Image.open(file_path) as img:
        img.convert('RGB').save(converted, 'JPEG')

    async with aiofiles.open(converted, 'rb') as f:                                        #changed with to async with and open to aiofiles.open  ..........maybe a problematic place!!!!!!!!!!!!!!!!!!!!!!!
        jpeg_bytes = f.read()

    save_to_io.write(jpeg_bytes)
    save_to_io.seek(0)
    os.remove(converted)

    return save_to_io.getvalue()


async def conv_to_pdf(file_path: str) -> bytes:
    pdf_filename = os.path.splitext(os.path.basename(file_path))[0] + '.pdf'

    # Создаем путь для сохранения PDF файла
    pdf_filepath = os.path.join(os.path.dirname(file_path), pdf_filename)

    # Конвертируем файл Word в PDF с помощью LibreOffice
    subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(file_path), file_path])

    # Читаем содержимое PDF файла в байтах и удаляем файл
    with open(pdf_filepath, 'rb') as f:
        pdf_bytes = f.read()
    os.remove(pdf_filepath)

    return pdf_bytes

# WORKING CODE HERE!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!


# async def conv_to_pdf(file_path: str, user_id) -> bytes:
#     save_to_io = io.BytesIO()
#     converted = fr'C:\PY\Python_learn\Minions_Bots\Converter_bot\tgbot_template_v3\save_dir\new_temp-{user_id}.pdf'
#     convert(file_path, converted)
#     with open(converted, 'rb') as f:
#         pdf_bytes = f.read()
#     save_to_io.write(pdf_bytes)
#     save_to_io.seek(0)
#     os.remove(converted)
#     return save_to_io.getvalue()






async def conv_to_docx(file_path: str, user_id) -> bytes:
    save_to_io = io.BytesIO()
    print('func started')
    converted = fr"C:\PY\Python_learn\Minions_Bots\Converter_bot\tgbot_template_v3\save_dir\new_temp-{user_id}.docx"
    await asyncio.to_thread(parse, file_path, converted)
    with open(converted, 'rb') as f:
        pdf_bytes = f.read()

    save_to_io.write(pdf_bytes)
    save_to_io.seek(0)
    os.remove(converted)
    print('docx created')
    return save_to_io.getvalue()


async def conv_to_txt(file_id, user_id) -> bytes:
    save_to_io = io.BytesIO()
    file = await bot.get_file(file_id)

    type_file = file.file_path.rsplit('.')[-1]
    logging.info('func started')
    logging.info(type_file)
    if type_file in ['png', 'jpeg', 'jpg']:
        async with aiofiles.open(file.file_path, 'rb') as f:               #WORKING CODE!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
            image_data = await f.read()

        img = Image.open(io.BytesIO(image_data))
        enhancer = ImageEnhance.Brightness(img)
        img = enhancer.enhance(1.5)
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(1.5)

        # Убрать шумы и преобразовать изображение в оттенки серого
        img = img.filter(ImageFilter.MedianFilter())
        gray = img.convert('L')

        # Применить фильтр Собеля для выделения границ
        img_arr = np.array(gray)
        sobelx = cv2.Sobel(img_arr, cv2.CV_64F, 1, 0, ksize=5)
        sobely = cv2.Sobel(img_arr, cv2.CV_64F, 0, 1, ksize=5)
        abs_grad_x = cv2.convertScaleAbs(sobelx)
        abs_grad_y = cv2.convertScaleAbs(sobely)
        sobel = cv2.addWeighted(abs_grad_x, 0.5, abs_grad_y, 0.5, 0)

        # Применить бинаризацию и распознать текст с помощью Tesseract OCR
        thresh = 200
        bw = gray.point(lambda x: 0 if x < thresh else 255, '1')
        text = pytesseract.image_to_string(bw, config='--psm 6 --oem 1')



        # image = Image.open(BytesIO(image_data))

        # text = pytesseract.image_to_string(image, config=custom_config)
        logging.info(text)
        if text != ' ':
            save_to_io.write(text.encode('utf-8'))
            save_to_io.seek(0)
            #os.remove(file_path)
            logging.info('txt created')
            return save_to_io.getvalue()
        else:
            logging.info(text)
            save_to_io.write('Found no text.\nPlease upload correct picture with text'.encode('utf-8'))
            save_to_io.seek(0)
            #os.remove(file_path)
            logging.info('txt created')
            return save_to_io.getvalue()
    elif 'docx' == type_file:
        text_doc = docx2txt.process(file.file_path)
        if text_doc:
            save_to_io.write(text_doc.encode('utf-8'))
            save_to_io.seek(0)
            #os.remove(file.file_path)
            logging.info('txt created')
            return save_to_io.getvalue()
        else:
            save_to_io.write('Found no text.\nPlease upload correct picture with text'.encode('utf-8'))
            save_to_io.seek(0)
            #os.remove(file.file_path)
            logging.info('txt created')
            return save_to_io.getvalue()


async def pptx_to_txt(file_path: str, user_id) -> bytes:
    save_to_io = io.BytesIO()
    print('func started')
    prs = Presentation(file_path)
    slide_texts = []
    for slide in prs.slides:
        slide_text = "\n".join([shape.text for shape in slide.shapes if shape.has_text_frame])
        slide_texts.append(slide_text)
    text = "\n\n".join(slide_texts)
    save_to_io.write(text.encode('utf-8'))
    save_to_io.seek(0)
    os.remove(file_path)
    print('txt created')
    return save_to_io.getvalue()


async def conv_to(file_path, file_id,  show_type_from, show_type_to, user_id) -> bytes:

    choose = {
        'docx': {
            'pdf': lambda: conv_to_pdf(file_path),
            'txt': lambda: conv_to_txt(file_id, user_id)
        },
        'pdf': {
            'docx': lambda: conv_to_docx(file_path, user_id)
        },
        'jpeg': {
            'txt': lambda: conv_to_txt(file_id, user_id)
        },

        'png':  {
            'txt': lambda: conv_to_txt(file_id, user_id),
            'jpeg': lambda: conv_to_jpeg(file_path, user_id)

        },
        'pptx': {
            'txt': lambda: pptx_to_txt(file_path, user_id)
        }

    }
    return await choose[show_type_from][show_type_to]()








#async def conv_xlsx_to_jpeg(file_path):
#     converted = r'C:\PY\Python_learn\Minions_Bots\Converter_bot\convert_directory\new_temp.jpeg'
#     print('func started')
#     # Открываем файл Excel с помощью openpyxl
#     async with aiofiles.open(file_path, mode='rb') as f:
#         wb = openpyxl.load_workbook(f, read_only=True)
#         # Получаем первый лист
#         ws = wb.active
#         # Устанавливаем размер изображения
#         img_width, img_height = 1200, 1600
#         # Создаем новое изображение
#         img = Image.new('RGB', (img_width, img_height), color='white')
#         # Получаем контекст изображения
#         img_draw = ImageDraw.Draw(img)
#         # Получаем данные из ячеек таблицы и записываем их в изображение
#         for row in ws.iter_rows(values_only=True):
#             for i, value in enumerate(row):
#                 x, y = i * 200, row[0] * 200
#                 # Добавляем значение ячейки в изображение
#                 img_draw.text((x, y), str(value), fill='black')
#         # Сохраняем изображение в формате JPEG
#         output_file_path = os.path.join(converted, os.path.splitext(os.path.basename(file_path))[0] + '.jpeg')
#         img.save(output_file_path, 'JPEG')
#         print('image saved')
#         return output_file_path


# async def conv_to_pdf(file_name, file, user_id):
#     text = docx2txt.process(file)
    # pdf_name = f"{file_name.strip('.')[0]}.pdf"
    # time_now = datetime.now()
    # pdf = FPDF()
    # pdf.add_page()
    # pdf.set_font("Arial", size=12)
    # pdf.cell(200, 10, txt=text, ln=1)
    # pdf.output('converted.pdf')
    # pdf_data = pdf.output(dest='S')
    #my_bytesio = io.BytesIO(b'{pdf_data}')
    #return pdf_data

# async def conv_to_pdf(file_name, file, user_id):
#     try:
#         text = docx2txt.process(file)
#         pdf_writer = PyPDF2.PdfWriter()
#         page_width = 595  # width of an A4 page in points
#         page_height = 842  # height of an A4 page in points
#         pdf_writer.add_page(PageObject.create_blank_page(pdf_writer, width=page_width, height=page_height))
#     except Exception as e:
#         pass

#print(conv_to_pdf('arc.doc', b"Hello", '112121'))


#save_file(pdf_name, bytes, user_id, time_now)

# async def conv_to_pdf(file_name, file, user_id):
#     text = docx2txt.process(file)
#     # pdf_name = f"{file_name.strip('.')[0]}.pdf"
#     # time_now = datetime.now()
#     pdf = FPDF()
#     pdf.add_page()
#     pdf.set_font("Arial", size=12)
#     pdf.cell(200, 10, txt=text, ln=1)
#     pdf_data = pdf.output(dest='S')
#     my_bytesio = io.BytesIO(pdf_data)
#     return my_bytesio

